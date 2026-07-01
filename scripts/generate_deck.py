#!/usr/bin/env python3
"""Generate a short presentation (PPTX) and PDF summarizing the Talent Intelligence System.

Writes `output/tis_deck.pptx` and `output/tis_deck.pdf` and includes a slide with
the top-10 ranked candidates read from `output/ranked_candidates.csv` if present.
"""
import csv
import os
import re
import textwrap
from pathlib import Path

OUT_DIR = Path("output")
OUT_DIR.mkdir(exist_ok=True)
PRESENTATION_SOURCE = Path("presentation/deck_outline.md")


def pdf_safe(text):
    return str(text).encode("latin-1", "replace").decode("latin-1")


def pdf_write_wrapped(pdf, text, width=92, line_height=6):
    safe_text = pdf_safe(text)
    wrapped = textwrap.wrap(
        safe_text,
        width=width,
        break_long_words=True,
        break_on_hyphens=True,
        drop_whitespace=False,
    )
    if not wrapped:
        pdf.ln(line_height)
        return
    for chunk in wrapped:
        pdf.multi_cell(0, line_height, chunk)


def read_top_candidates(csv_path, n=10):
    rows = []
    if not csv_path.exists():
        return rows
    with csv_path.open(newline='', encoding='utf-8') as fh:
        reader = csv.DictReader(fh)
        for i, r in enumerate(reader):
            rows.append(r)
            if i + 1 >= n:
                break
    return rows


def read_presentation_source():
    for candidate in [PRESENTATION_SOURCE, Path("README.md")]:
        if candidate.exists():
            return candidate.read_text(encoding="utf-8")
    return ""


def parse_slide_sections(text):
    slides = []
    current = None
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            if current is not None:
                current["body"].append("")
            continue
        if line.startswith("## Slide "):
            if current is not None:
                slides.append(current)
            title = line.split("—", 1)[-1].strip() if "—" in line else line.replace("## ", "").strip()
            current = {"title": title, "body": []}
        elif current is not None:
            current["body"].append(line)
    if current is not None:
        slides.append(current)
    return slides


def split_bullets(lines):
    bullets = []
    for line in lines:
        if line.startswith("- "):
            bullets.append(line[2:].strip())
    return bullets


def build_presentation_fallback(top_candidates, readme_text):
    # Lightweight fallback when python-pptx can't be used; write a short text file
    txt = OUT_DIR / "tis_deck.txt"
    with txt.open("w", encoding="utf-8") as fh:
        fh.write("Talent Intelligence System — Summary\n\n")
        fh.write(readme_text[:8000])
        fh.write("\n\nTop candidates:\n")
        for r in top_candidates:
            reason = r.get('reason', r.get('reasoning', ''))
            fh.write(f"{r.get('candidate_id','?')}: {r.get('score','?')} — {reason}\n")
    return txt


def build_pdf(top_candidates, readme_text, out_pdf_path):
    try:
        from reportlab.lib.pagesizes import LETTER
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer
    except Exception:
        from fpdf import FPDF

        pdf = FPDF()
        pdf.set_auto_page_break(True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, pdf_safe("Talent Intelligence System - Summary"), ln=True)
        pdf.ln(4)
        pdf.set_font("Arial", size=11)

        slides = parse_slide_sections(readme_text)
        for slide in slides:
            pdf.set_font("Arial", "B", 13)
            pdf_write_wrapped(pdf, slide["title"], width=88, line_height=7)
            pdf.set_font("Arial", size=11)
            bullets = split_bullets(slide["body"])
            if bullets:
                for bullet in bullets:
                    pdf_write_wrapped(pdf, f"- {bullet}", width=88, line_height=6)
            else:
                for line in slide["body"]:
                    if line.startswith("#"):
                        continue
                    pdf_write_wrapped(pdf, line, width=88, line_height=6)
            pdf.ln(2)
        pdf.add_page()
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 8, pdf_safe("Top Candidates (from output/ranked_candidates.csv)"), ln=True)
        pdf.ln(2)
        pdf.set_font("Arial", size=11)
        if not top_candidates:
            pdf_write_wrapped(pdf, "No ranked candidates file found. Run python run.py --dry-run first.", width=88, line_height=6)
        else:
            for r in top_candidates:
                cid = r.get("candidate_id", "?")
                score = r.get("score", "?")
                reason = r.get("reason", r.get("reasoning", ""))
                pdf_write_wrapped(pdf, f"{cid} - score: {score} - {reason}", width=88, line_height=6)

        pdf.output(str(out_pdf_path))
        return

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "DeckTitle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=22,
        leading=26,
        spaceAfter=12,
    )
    heading_style = ParagraphStyle(
        "DeckHeading",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=15,
        leading=18,
        spaceBefore=6,
        spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "DeckBody",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10.5,
        leading=13,
        spaceAfter=3,
    )
    bullet_style = ParagraphStyle(
        "DeckBullet",
        parent=body_style,
        leftIndent=12,
        bulletIndent=0,
        spaceAfter=2,
    )

    doc = SimpleDocTemplate(
        str(out_pdf_path),
        pagesize=LETTER,
        rightMargin=0.7 * inch,
        leftMargin=0.7 * inch,
        topMargin=0.65 * inch,
        bottomMargin=0.65 * inch,
    )

    story = [Paragraph("Talent Intelligence System - Summary", title_style), Spacer(1, 0.12 * inch)]
    slides = parse_slide_sections(readme_text)
    for slide in slides:
        story.append(Paragraph(slide["title"], heading_style))
        bullets = split_bullets(slide["body"])
        if bullets:
            for bullet in bullets:
                story.append(Paragraph(bullet, bullet_style, bulletText="-"))
        else:
            for line in slide["body"]:
                if line.startswith("#"):
                    continue
                story.append(Paragraph(line, body_style))
        story.append(Spacer(1, 0.08 * inch))

    story.append(PageBreak())
    story.append(Paragraph("Top Candidates (from output/ranked_candidates.csv)", heading_style))
    if not top_candidates:
        story.append(Paragraph("No ranked candidates file found. Run python run.py --dry-run first.", body_style))
    else:
        for r in top_candidates:
            cid = r.get("candidate_id", "?")
            score = r.get("score", "?")
            reason = r.get("reason", r.get("reasoning", ""))
            story.append(Paragraph(f"{cid} - score: {score} - {reason}", body_style))

    doc.build(story)


def build_pptx_full(top_candidates, readme_text, out_pptx_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slides = parse_slide_sections(readme_text)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Talent Intelligence System"
    try:
        subtitle = slide.placeholders[1]
        subtitle.text = "AI-powered candidate ranking, explained"
    except Exception:
        pass

    for slide_info in slides:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide_info["title"]
        tf = s.shapes.placeholders[1].text_frame
        body_lines = slide_info["body"]
        bullets = split_bullets(body_lines)
        if bullets:
            tf.text = bullets[0]
            for bullet in bullets[1:]:
                p = tf.add_paragraph()
                p.level = 0
                p.text = bullet
        else:
            tf.text = body_lines[0] if body_lines else ""
            for line in body_lines[1:]:
                if line.startswith("#"):
                    continue
                p = tf.add_paragraph()
                p.level = 0
                p.text = line

    # Top candidates slide
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Top Candidates (sample)"
    tf3 = s3.shapes.placeholders[1].text_frame
    if not top_candidates:
        tf3.text = "No ranked candidates found. Run pipeline to generate output/ranked_candidates.csv"
    else:
        tf3.text = "Top results:"
        for r in top_candidates:
            p = tf3.add_paragraph()
            cid = r.get("candidate_id", "?")
            score = r.get("score", "?")
            reason = r.get("reason", r.get("reasoning", ""))
            p.level = 0
            p.text = f"{cid} — {score} — {reason[:120]}"

    prs.save(str(out_pptx_path))


def main():
    csv_path = Path("output/ranked_candidates.csv")
    top = read_top_candidates(csv_path, n=10)
    readme_text = read_presentation_source()
    if not readme_text:
        readme_text = "(presentation source not found)"

    out_pptx = OUT_DIR / "tis_deck.pptx"
    out_pdf = OUT_DIR / "tis_deck.pdf"

    # Try to write PPTX and PDF; if python-pptx not available, write a fallback text and still create PDF if possible
    try:
        build_pptx_full(top, readme_text, out_pptx)
    except Exception:
        try:
            build_presentation_fallback(top, readme_text)
        except Exception:
            pass

    try:
        build_pdf(top, readme_text, out_pdf)
    except Exception:
        # if PDF build fails, write fallback text
        fb = build_presentation_fallback(top, readme_text)
        print(f"Wrote fallback text summary to {fb}")

    print(f"Wrote: {out_pptx} and {out_pdf} (if creation succeeded)")


if __name__ == "__main__":
    main()
